"""Script para generar la PPT de presentacion del sistema experto."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

C_BG     = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0xE9, 0x4F, 0x37)
C_LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)
C_MID    = RGBColor(0x16, 0x21, 0x3E)
C_SUB    = RGBColor(0xAA, 0xAA, 0xCC)
C_DIM    = RGBColor(0x88, 0x88, 0xAA)
C_CODE   = RGBColor(0x7E, 0xE8, 0x7E)
C_CODEBG = RGBColor(0x0D, 0x11, 0x17)

blank = prs.slide_layouts[6]


def rect(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, l, t, w, h, size=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def footer(slide):
    txt(slide, "Triage Mecanico — Sistema Experto | AD2 2026",
        0.4, 7.1, 12, 0.4, size=10, color=C_DIM, align=PP_ALIGN.LEFT)


def base_bg(slide, bar_accent=True):
    rect(slide, 0, 0, 13.33, 7.5, fill=C_BG)
    if bar_accent:
        rect(slide, 0, 0, 13.33, 0.08, fill=C_ACCENT)
    rect(slide, 0, 7.3, 13.33, 0.2, fill=C_MID)
    footer(slide)


def make_title_slide(title, subtitle="", blk=None):
    s = prs.slides.add_slide(blank)
    base_bg(s)
    if blk:
        rect(s, 0.4, 1.5, 0.7, 0.7, fill=C_ACCENT)
        txt(s, str(blk), 0.4, 1.5, 0.7, 0.7, size=22, bold=True,
            color=C_LIGHT, align=PP_ALIGN.CENTER)
    lft = 1.3 if blk else 0.6
    txt(s, title, lft, 1.4, 11, 1.8, size=36, bold=True,
        color=C_LIGHT, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(s, subtitle, lft, 3.3, 11, 1.2, size=20, color=C_SUB,
            align=PP_ALIGN.LEFT)
    return s


def make_content_slide(title, bullets, blk=None, code=None):
    s = prs.slides.add_slide(blank)
    base_bg(s)
    rect(s, 0, 0.08, 0.08, 7.22, fill=C_MID)
    if blk:
        rect(s, 0.2, 0.2, 0.55, 0.55, fill=C_ACCENT)
        txt(s, str(blk), 0.2, 0.2, 0.55, 0.55, size=16, bold=True,
            color=C_LIGHT, align=PP_ALIGN.CENTER)
        txt(s, title, 0.9, 0.18, 11.5, 0.65, size=22, bold=True,
            color=C_LIGHT, align=PP_ALIGN.LEFT)
    else:
        txt(s, title, 0.3, 0.18, 12.5, 0.65, size=22, bold=True,
            color=C_LIGHT, align=PP_ALIGN.LEFT)
    rect(s, 0.3, 0.95, 12.8, 0.04, fill=C_ACCENT)

    content_w = 7.8 if code else 12.8
    top = 1.15
    for b in bullets:
        is_sub = b.startswith("  ")
        btext = b.strip()
        bsize = 14 if is_sub else 17
        bleft = 0.8 if is_sub else 0.5
        indent = "       " if is_sub else "   "
        color = RGBColor(0xAA, 0xAA, 0xCC) if is_sub else RGBColor(0xDD, 0xDD, 0xEE)
        txt(s, indent + btext, bleft, top, content_w, 0.45, size=bsize, color=color)
        top += 0.42 if is_sub else 0.47

    if code:
        rect(s, 8.8, 1.0, 4.2, 5.8, fill=C_CODEBG)
        txt(s, code, 8.95, 1.1, 4.0, 5.6, size=9, color=C_CODE)
    return s


# ── PORTADA ──────────────────────────────────────────────────────────────────
s0 = prs.slides.add_slide(blank)
base_bg(s0)
rect(s0, 5.5, 1.5, 2.3, 0.08, fill=C_ACCENT)
txt(s0, "TRIAGE MECANICO", 1.5, 1.8, 10, 1.5, size=52, bold=True,
    color=C_LIGHT, align=PP_ALIGN.CENTER)
txt(s0, "Sistema Experto de Diagnostico Automotriz", 1.5, 3.5, 10, 0.8,
    size=26, color=C_SUB, align=PP_ALIGN.CENTER)
rect(s0, 5.5, 4.4, 2.3, 0.08, fill=C_ACCENT)
txt(s0, "Analisis de Datos II — Sistemas Expertos | 2026", 1.5, 4.7, 10, 0.5,
    size=16, color=C_DIM, align=PP_ALIGN.CENTER)

# ── BLOQUE 1: Introduccion ───────────────────────────────────────────────────
make_title_slide("Introduccion y Problema del Dominio",
                 "Por que un Sistema Experto para diagnostico mecanico?", 1)

make_content_slide("Por que IA Simbolica?", [
    "El diagnostico mecanico requiere experiencia experta",
    "  La mayoria de conductores no sabe si pueden seguir manejando",
    "",
    "Feigenbaum (1982): un SE usa conocimiento y procedimientos",
    "  de inferencia para resolver problemas que requieren experticia humana",
    "",
    "Ventajas sobre Deep Learning para este dominio:",
    "  Interpretable: cada diagnostico tiene una regla identificable",
    "  Explicable: tabla de trazabilidad muestra el razonamiento",
    "  Mantenible: agregar fallas = escribir nuevas reglas",
    "  No requiere datos etiquetados — el conocimiento viene del experto",
], blk=1)

make_content_slide("Dominio y Alcance", [
    "6 sistemas del vehiculo cubiertos:",
    "  Frenos | Motor | Refrigeracion | Transmision | Suspension | Electrico",
    "",
    "10 tipos de sintomas de entrada:",
    "  Ruidos: chirrido, golpeteo, zumbido, crujido",
    "  Humo: azul, blanco, negro",
    "  Temperatura del motor: normal / alta / muy alta",
    "  Luces del tablero: aceite, temperatura, bateria, check engine",
    "  Vibracion: volante / pedal / todo el auto",
    "  Comportamiento del freno y liquido en el piso",
    "",
    "28 reglas — 4 niveles de urgencia: critica / alta / moderada / baja",
], blk=1)

# ── BLOQUE 2: Base de Conocimiento ───────────────────────────────────────────
make_title_slide("Base de Conocimiento",
                 "Estructura del RBES — memoria a largo y corto plazo", 2)

make_content_slide("Componentes de la Base de Conocimiento", [
    "BASE DE REGLAS: 28 reglas IF-THEN organizadas en 6 modulos",
    "  SI temperatura=muy_alta Y liquido_piso=refrigerante",
    "  Y perdida_potencia=True ENTONCES falla_multiple (critica)",
    "",
    "BASE DE HECHOS GENERALES (memoria a largo plazo):",
    "  Catalogo _RESULTADOS en motor.py",
    "  Contiene: diagnostico, urgencia, seguro_manejar, accion recomendada",
    "",
    "HECHOS DEL CASO (working memory / memoria de corto plazo):",
    "  Objeto _Sintoma(...) declarado al inicio de cada consulta",
    "  Los sintomas del conductor para esa sesion especifica",
    "",
    "BASE DE PREGUNTAS:",
    "  La interfaz Streamlit funciona como modulo de preguntas",
], blk=2,
code=(
    "# Regla IF-THEN en Python\n"
    "@Rule(\n"
    "  _Sintoma(\n"
    "    temperatura=TEMP_MUY_ALTA,\n"
    "    perdida_potencia=True,\n"
    "    liquido_piso=LIQUIDO_REFRIGERANTE\n"
    "  ),\n"
    "  NOT(_Diagnostico()),\n"
    "  salience=97,\n"
    ")\n"
    "def refrigeracion_falla_multiple(self):\n"
    "  self.declare(\n"
    "    _Diagnostico(\n"
    "      nombre=\"refrigeracion\n"
    "        .falla_multiple\"\n"
    "    )\n"
    "  )"
))

# ── BLOQUE 3: Motor de Inferencia ─────────────────────────────────────────────
make_title_slide("Motor de Inferencia",
                 "Forward Chaining — Ciclo Match-Resolve-Act", 3)

make_content_slide("Ciclo Match-Resolve-Act", [
    "1. MATCH (Coincidencia de patrones):",
    "  Compara antecedentes de cada @Rule con los hechos del caso",
    "  En Experta: cada regla evaluada contra el _Sintoma declarado",
    "",
    "2. RESOLVE (Resolucion de conflictos — estrategia SALIENCE):",
    "  Importancia: la regla con mayor puntaje de prioridad gana",
    "  Pastillas y discos: salience=100 (5 condiciones especificas)",
    "  Recalentamiento critico: salience=93 (condicion termica)",
    "  Desbalanceo: salience=52 (condicion generica)",
    "",
    "3. ACT (Activacion de consecuencia):",
    "  La regla ganadora declara _Diagnostico(nombre=...)",
    "  NOT(_Diagnostico()) garantiza UNA sola regla activa",
], blk=3,
code=(
    "# Regla alta prioridad\n"
    "@Rule(\n"
    "  _Sintoma(\n"
    "    temperatura=TEMP_MUY_ALTA,\n"
    "    perdida_potencia=True,\n"
    "    liquido_piso=LIQUIDO_REFRIGERANTE\n"
    "  ),\n"
    "  NOT(_Diagnostico()),\n"
    "  salience=97,  # alta\n"
    ")\n"
    "def falla_multiple(self):\n"
    "  self.declare(...)\n\n"
    "# Regla menos especifica\n"
    "@Rule(\n"
    "  _Sintoma(temperatura=TEMP_MUY_ALTA),\n"
    "  NOT(_Diagnostico()),\n"
    "  salience=93,  # menor\n"
    ")\n"
    "def recalentamiento(self):\n"
    "  self.declare(...)"
))

make_content_slide("Experta como Shell del SE", [
    "Experta: alternativa Python a CLIPS (visto en Clase 5)",
    "  Motor de inferencia tipo forward chaining",
    "  Basado en PyKnow (2018), compatible con Python 3",
    "",
    "Componentes de Experta que usamos:",
    "  Fact: clase base para _Sintoma y _Diagnostico",
    "  KnowledgeEngine: motor con ciclo reset-declare-run",
    "  Rule: decorador que define condicion + accion",
    "  NOT: condicion negativa (no existe hecho de tipo X)",
    "  salience: puntaje de prioridad para resolucion de conflictos",
    "",
    "Parche de compatibilidad Python 3.10+:",
    "  collections.Mapping removido — restaurado manualmente",
    "  antes de importar experta para evitar ImportError",
], blk=3)

# ── BLOQUE 4: Innovaciones ────────────────────────────────────────────────────
make_title_slide("Innovaciones del Sistema",
                 "Sintomas opcionales + Fuzzy Scoring + Reglas multiples", 4)

make_content_slide("Sintomas Visuales Opcionales", [
    "PROBLEMA: RBES clasico falla si falta cualquier condicion",
    "  La luz de temperatura requiere que el foco funcione",
    "  Si el foco esta quemado — la regla no dispara aunque el motor se queme",
    "",
    "SOLUCION: campos opcionales en la evaluacion de reglas",
    "  _CAMPOS_OPCIONALES = {'luces'}",
    "  Si luces='ninguna' → condicion EXCLUIDA del chequeo",
    "  La regla puede disparar solo con sintomas fisicos",
    "",
    "Motor y Refrigeracion: reglas sin condicion de luz",
    "  Detectan recalentamiento con solo temperatura=muy_alta",
    "Electrico: si requiere luz (es la premisa principal)",
    "  bateria_alternador solo dispara si luz de bateria esta encendida",
], blk=4,
code=(
    "# Antes (bloqueante):\n"
    "@Rule(\n"
    "  _Sintoma(\n"
    "    temperatura=TEMP_MUY_ALTA,\n"
    "    luces=LUZ_TEMPERATURA  # BLOQUEA\n"
    "  ),\n"
    "  salience=93,\n"
    ")\n\n"
    "# Ahora (opcional):\n"
    "@Rule(\n"
    "  _Sintoma(temperatura=TEMP_MUY_ALTA),\n"
    "  # luces no en Rule\n"
    "  salience=93,\n"
    ")\n\n"
    "# En _encontrar_compatibles:\n"
    "if campo in _CAMPOS_OPCIONALES\\\n"
    "   and sintomas[campo] == 'ninguna':\n"
    "    continue  # no bloquea"
))

make_content_slide("Fuzzy Scoring y Reglas Multiples", [
    "FUZZY SCORING — fallback cuando el engine no dispara exactamente:",
    "  Inspirado en los Fuzzy ES de la Clase 5",
    "  Para cada regla: confianza = condiciones_cumplidas / activas",
    "  Umbral: si alguna regla supera 70%, es candidata",
    "  Gana la de mayor urgencia; en empate, la de mayor confianza",
    "  Los campos opcionales se excluyen del denominador si no reportados",
    "",
    "TODAS LAS REGLAS COMPATIBLES — innovacion de esta version:",
    "  _encontrar_compatibles(sintomas) evalua TODAS las reglas",
    "  Resultado: diagnostico = el mas urgente (salience maxima)",
    "  Tabla reglas_disparadas = TODAS las compatibles ordenadas",
    "  Ejemplo: frenos.liquido_bajo + refrigeracion.recalentamiento_critico",
    "  El conductor ve ambos problemas aunque el diagnostico sea uno solo",
], blk=4)

# ── BLOQUE 5: Sistema de Explicacion e Interfaz ───────────────────────────────
make_title_slide("Sistema de Explicacion e Interfaz",
                 "Trazabilidad del razonamiento + Streamlit", 5)

make_content_slide("Sistema de Explicacion", [
    "Componente clave de todo SE: el usuario debe entender el POR QUE",
    "  En aplicaciones criticas (salud, seguridad) la explicabilidad es esencial",
    "",
    "Tabla 'Reglas que se dispararon':",
    "  Regla: identificador (ej: refrigeracion.recalentamiento_critico)",
    "  Sistema: modulo afectado (Refrigeracion, Frenos, Motor...)",
    "  Diagnostico asociado: descripcion de la falla",
    "  Urgencia: critica/alta/moderada/baja con icono de color",
    "",
    "Si hay multiples reglas compatibles, todas aparecen",
    "  ordenadas por urgencia descendente (mas critica primero)",
    "",
    "Esto permite al mecanico revisar TODAS las areas afectadas",
    "  no solo la de mayor prioridad — vision completa del problema",
], blk=5)

make_content_slide("Interfaz con Streamlit", [
    "Framework Python para apps web interactivas (Clase 5)",
    "  pip install streamlit → streamlit run app.py",
    "",
    "Componentes de la UI:",
    "  Sidebar: checkboxes por categoria + campos condicionales",
    "  Badge de urgencia: color semaforo (rojo/naranja/amarillo/verde)",
    "  Gauge Plotly: nivel de urgencia en escala 0-100",
    "  4 Cards: sistema afectado, diagnostico, seguridad, accion",
    "  Tabla de reglas: DataFrame con columnas de contexto",
    "  5 Casos demo: escenarios precargados para demostracion rapida",
    "",
    "Flujo del usuario:",
    "  1. Seleccionar categorias de sintomas en sidebar",
    "  2. Completar campos especificos de cada categoria",
    "  3. Presionar Diagnosticar → ver resultado con trazabilidad",
], blk=5)

# ── BLOQUE 6: Resultados y Conclusiones ───────────────────────────────────────
make_title_slide("Resultados y Conclusiones",
                 "Arquitectura final, QA y aprendizajes del proyecto", 6)

make_content_slide("Arquitectura Final", [
    "Sintomas del conductor (Streamlit UI)",
    "  |",
    "  v  src/motor.py (orquestador)",
    "  [1] TriageEngine (Experta) — Forward Chaining + Salience",
    "       Regla ganadora = diagnostico principal",
    "  [2] _encontrar_compatibles() — todas las reglas compatibles",
    "       Usadas para la tabla de trazabilidad",
    "  [3] Fallback: _score_fuzzy() — si no hay match exacto",
    "       Logica difusa con umbral 70%",
    "  |",
    "  v  app.py (Streamlit)",
    "  Badge urgencia + Gauge + 4 Cards + Tabla reglas disparadas",
], blk=6,
code=(
    "# Flujo diagnosticar()\n"
    "engine = TriageEngine()\n"
    "engine.reset()\n"
    "engine.declare(_Sintoma(**s))\n"
    "engine.run()\n\n"
    "# Regla ganadora\n"
    "ganador = None\n"
    "for fact in engine.facts.values():\n"
    "  if isinstance(fact, _Diagnostico):\n"
    "    ganador = fact['nombre']\n\n"
    "# Todas las compatibles\n"
    "todas = _encontrar_compatibles(s)\n\n"
    "if ganador:\n"
    "  return _construir_resultado(\n"
    "    ganador, todas\n"
    "  )\n\n"
    "# Fallback fuzzy\n"
    "return _score_fuzzy(s)"
))

make_content_slide("Resultados QA y Conclusiones", [
    "RESULTADOS QA:",
    "  30 casos de prueba en data/casos_prueba.csv",
    "  28/30 correctos — 93% de tasa de acierto",
    "  Los 2 casos restantes: inconsistencias del dataset",
    "  (mismos sintomas mapeados a diagnosticos distintos)",
    "",
    "CONCLUSIONES:",
    "  RBES ideal cuando el conocimiento es explícito y la explicabilidad importa",
    "  Extension fuzzy maneja imprecision del diagnostico de campo",
    "  Sintomas opcionales modelan informacion incompleta del conductor",
    "  Todas las reglas compatibles dan vision completa del problema",
    "",
    "TRABAJO FUTURO:",
    "  Mas sintomas (ruidos especificos, codigos OBD)",
    "  Neural-ES hibrido para aprendizaje inductivo",
], blk=6)

# ── SLIDE FINAL ───────────────────────────────────────────────────────────────
sf = prs.slides.add_slide(blank)
base_bg(sf)
txt(sf, "Gracias", 1.5, 2.0, 10, 1.8, size=60, bold=True,
    color=C_LIGHT, align=PP_ALIGN.CENTER)
txt(sf, "Preguntas?", 1.5, 3.9, 10, 0.9, size=30,
    color=C_SUB, align=PP_ALIGN.CENTER)
txt(sf, "Analisis de Datos II — Sistemas Expertos | AD2 2026",
    1.5, 5.0, 10, 0.5, size=16, color=C_DIM, align=PP_ALIGN.CENTER)

prs.save("presentaciones-profesor/presentacion_SE.pptx")
print("PPT generada correctamente.")
